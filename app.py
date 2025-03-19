import streamlit as st
import pandas as pd

# Configuración básica de la aplicación
st.set_page_config(page_title="Generador de Certificados", layout="centered")

st.title("🎓 Generador de Certificados")

st.write("Bienvenido al generador de certificados. Ingrese la información requerida para generar su diploma.")

# URL pública de Google Sheets en formato CSV
url = "https://docs.google.com/spreadsheets/d/1XSzJ_cZWr7co6c_86CCzfWNgEKwxB8Wn5NBt4PpNAUc/gviz/tq?tqx=out:csv"

# Cargar la hoja de cálculo como un DataFrame
@st.cache_data
def load_data():
    return pd.read_csv(url)

df_cursos = load_data()

# Cargar nuevamente los datos sin caché para verificar actualización
df_cursos = pd.read_csv(url)  # Sin @st.cache_data

st.write("Columnas disponibles en la hoja de cálculo:", df_cursos.columns.tolist())

# Mostrar los primeros registros para verificar
st.subheader("Lista de cursos disponibles")
st.dataframe(df_cursos[["Código", "Nombre del Curso o Diplomado", "Cohorte", "Fecha", "Duración","Validación"]])

# Seleccionar el curso desde un selectbox
curso_seleccionado = st.selectbox("📚 Seleccione un curso o diplomado", df_cursos["Nombre del Curso o Diplomado"].unique())

# Obtener el código del curso seleccionado
codigo_curso = df_cursos[df_cursos["Nombre del Curso o Diplomado"] == curso_seleccionado]["Código"].values[0]

st.write(f"Has seleccionado el curso: **{curso_seleccionado}**")
st.write(f"Código del curso: `{codigo_curso}`")

# URL pública de la hoja de cálculo de estudiantes en formato CSV
url_estudiantes = "https://docs.google.com/spreadsheets/d/1prUt0i0EWolsX_LuGl_yKzXPUWmy6CzCxi28zued5BA/gviz/tq?tqx=out:csv"

# Cargar los datos de los estudiantes
@st.cache_data
def load_students():
    return pd.read_csv(url_estudiantes)

df_estudiantes = load_students()

# Definir variables globales para que estén disponibles en toda la app
estudiante = pd.DataFrame()
nombre_estudiante = ""
documento_estudiante = ""

# Inicializar valores en session_state si no existen
if "nombre_estudiante" not in st.session_state:
    st.session_state.nombre_estudiante = ""
    st.session_state.documento_estudiante = ""
    st.session_state.validado = False

# Campo de entrada para la contraseña del estudiante
password_input = st.text_input("🔑 Ingrese su contraseña", type="password")

# Filtrar la hoja de estudiantes por el código del curso seleccionado
df_curso_estudiantes = df_estudiantes[df_estudiantes["Código"] == codigo_curso]

# Botón para validar
if st.button("Validar contraseña"):
    if password_input:
        estudiante = df_curso_estudiantes[df_curso_estudiantes["Contraseña"] == password_input]

        if not estudiante.empty:
            # Guardar los datos en session_state para que se conserven entre interacciones
            st.session_state.nombre_estudiante = estudiante["Nombre"].values[0]
            st.session_state.documento_estudiante = estudiante["Documento"].values[0]
            st.session_state.validado = True

            st.success(f"✅ Acceso concedido: {st.session_state.nombre_estudiante}")
            st.write(f"📄 Documento: `{st.session_state.documento_estudiante}`")
        else:
            st.session_state.validado = False
            st.error("❌ Contraseña incorrecta o estudiante no registrado en este curso.")
    else:
        st.warning("⚠️ Por favor, ingrese su contraseña.")

import qrcode
from io import BytesIO

# Obtener el enlace de validación para el curso seleccionado
url_validacion = df_cursos[df_cursos["Código"] == codigo_curso]["Validación"].values[0]

if "http" in url_validacion:  # Verificar que sea un enlace válido
    # Generar el código QR
    qr = qrcode.make(url_validacion)
    qr_img = BytesIO()
    qr.save(qr_img, format="PNG")
    qr_img.seek(0)

    # Mostrar el código QR en Streamlit
    st.subheader("Código QR de validación")
    st.image(qr_img, caption="Escanéalo para verificar tu certificado", use_container_width=True)
else:
    st.warning("⚠️ Este curso no tiene un enlace de validación asignado.")

import requests
from pptx import Presentation

# URL pública del archivo PPTX en GitHub
url_plantilla = "https://raw.githubusercontent.com/ciec-uptc/Certificados/main/Plantilla%20base.pptx"

def load_template():
    response = requests.get(url_plantilla)
    if response.status_code == 200:
        with open("Plantilla_base.pptx", "wb") as f:
            f.write(response.content)
        return Presentation("Plantilla_base.pptx")
    else:
        st.error("❌ No se pudo descargar la plantilla del certificado.")
        return None

# Cargar la plantilla sin almacenamiento en caché
plantilla_pptx = load_template()

import requests
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

# URL de la plantilla base en GitHub (raw)
url_plantilla = "https://github.com/ciec-uptc/Certificados/blob/main/Plantilla%20base.pptx?raw=true"

# Descargar la plantilla base desde GitHub
@st.cache_data
def cargar_plantilla():
    response = requests.get(url_plantilla)
    if response.status_code == 200:
        return BytesIO(response.content)
    else:
        st.error("❌ Error al cargar la plantilla de PowerPoint.")
        return None

plantilla_stream = cargar_plantilla()

from pptx.dml.color import RGBColor

#Generar certificado
def generar_certificado(nombre, documento, curso, duracion, fecha, qr_img):
    if plantilla_stream:
        prs = Presentation(plantilla_stream)  # Cargar la plantilla en memoria

        # Asegurar que los valores no sean NaN ni None
        nombre = str(nombre) if pd.notna(nombre) else ""
        documento = str(documento) if pd.notna(documento) else ""
        curso = str(curso) if pd.notna(curso) else ""
        duracion = str(duracion) if pd.notna(duracion) else ""
        fecha = str(fecha) if pd.notna(fecha) else ""

        # Modificar los textos sin perder formato
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text:
                    text = shape.text_frame.text.strip()

                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:  # Mantiene el formato original
                            if "Nombres y Apellidos" in text:
                                run.text = nombre
                            elif "Documento" in text:
                                run.text = documento
                            elif "Título" in text:
                                run.text = curso
                            elif "Dur" in text:
                                run.text = duracion
                            elif "Fecha" in text:
                                run.text = fecha

        # Insertar el código QR reemplazando el cuadro de texto "QR AQUÍ"
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and "QR AQUÍ" in shape.text_frame.text:
                    left, top, width, height = shape.left, shape.top, shape.width, shape.height

                    # Eliminar el cuadro de texto original
                    slide.shapes._spTree.remove(shape._element)

                    # Guardar el QR como imagen
                    qr_stream = BytesIO()
                    qr_img.save(qr_stream, format="PNG")
                    qr_stream.seek(0)

                    # Insertar el QR en la misma posición
                    slide.shapes.add_picture(qr_stream, left, top, width, height)
                    break  # Detener la búsqueda después de insertar el QR

        # Guardar el certificado como un archivo en memoria
        certificado_stream = BytesIO()
        prs.save(certificado_stream)
        certificado_stream.seek(0)
        
        return certificado_stream
    else:
        st.error("❌ No se pudo generar el certificado.")
        return None

# Botón para generar el certificado y permitir la descarga
if st.session_state.validado:
    if st.button("📜 Generar Certificado"):
        certificado_stream = generar_certificado(
            st.session_state.nombre_estudiante,
            st.session_state.documento_estudiante,
            curso_seleccionado,
            df_cursos[df_cursos["Código"] == codigo_curso]["Duración"].values[0],
            df_cursos[df_cursos["Código"] == codigo_curso]["Fecha"].values[0],
            qr
        )

        if certificado_stream:
            st.success("✅ Certificado generado con éxito.")
            
            # Botón de descarga
            st.download_button(
                label="⬇️ Descargar Certificado",
                data=certificado_stream,
                file_name=f"Certificado_{st.session_state.nombre_estudiante}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

import streamlit as st
from pptx import Presentation
import io
import cv2
import numpy as np
from PIL import Image

def convertir_pptx_a_video(certificado_stream):
    """Convierte la diapositiva PPTX en un video MP4 de alta calidad sin perder diseño."""

    # Guardar el archivo PPTX temporalmente
    pptx_path = "certificado_temporal.pptx"
    with open(pptx_path, "wb") as f:
        f.write(certificado_stream.getbuffer())

    # Cargar la presentación
    prs = Presentation(pptx_path)
    slide = prs.slides[0]  # Primera y única diapositiva

    # Dimensiones en píxeles (PowerPoint usa puntos de 1/72 pulgadas)
    width_px = int(prs.slide_width.inches * 96)
    height_px = int(prs.slide_height.inches * 96)

    # Crear una imagen en blanco
    img = Image.new("RGB", (width_px, height_px), "white")

    # Renderizar imágenes y texto de la diapositiva
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            left = int(shape.left.inches * 96)
            top = int(shape.top.inches * 96)
            draw = ImageDraw.Draw(img)
            draw.text((left, top), text, fill="black")

        if shape.shape_type == 13:  # Si es una imagen
            img_stream = io.BytesIO(shape.image.blob)
            image_pil = Image.open(img_stream).convert("RGBA")
            img.paste(image_pil, (shape.left, shape.top), image_pil)

    # Convertir la imagen a un array numpy para OpenCV
    frame = np.array(img)

    # Crear el video con OpenCV
    video_path = "certificado.mp4"
    fourcc = cv2.VideoWriter_fourcc(*"mp4v")  # Codec MP4
    fps = 1  # 1 frame por segundo
    duration = 5  # 5 segundos de duración
    out = cv2.VideoWriter(video_path, fourcc, fps, (width_px, height_px))

    for _ in range(fps * duration):  # Agregar múltiples frames para duración
        out.write(frame)

    out.release()

    # Leer el video generado en memoria
    video_buffer = io.BytesIO()
    with open(video_path, "rb") as f:
        video_buffer.write(f.read())
    video_buffer.seek(0)

    return video_buffer

# Generar el certificado en PPTX
if st.session_state.validado:
    certificado_stream = generar_certificado(
        st.session_state.nombre_estudiante,
        st.session_state.documento_estudiante,
        curso_seleccionado,
        df_cursos[df_cursos["Código"] == codigo_curso]["Duración"].values[0],
        df_cursos[df_cursos["Código"] == codigo_curso]["Fecha"].values[0],
        qr
    )

    if certificado_stream:
        st.success("✅ Certificado generado con éxito.")

        # Convertir PPTX a MP4 manteniendo TODO el diseño
        certificado_video = convertir_pptx_a_video(certificado_stream)

        # Botón de descarga en Streamlit
        st.download_button(
            label="🎥 Descargar Certificado en Video MP4",
            data=certificado_video,
            file_name=f"Certificado_{st.session_state.nombre_estudiante}.mp4",
            mime="video/mp4"
        )
